"""Generate PowerPoint and PDF deliverables for PEINP Parks Meteo Optimization."""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from fpdf import FPDF

OUT_DIR = Path("outputs")
OUT_DIR.mkdir(exist_ok=True)
DATA_DIR = Path("data/processed")

# -- Shared content ----------------------------------------------------------

TITLE = "Parks Meteo Optimization"
SUBTITLE = "PEINP Meteorological Data Pipeline - Final Report"
DATE = "April 2026"

SLIDES = [
    {
        "title": "Project Overview",
        "bullets": [
            "End-to-end pipeline for meteorological data analysis in Prince Edward Island National Park",
            "5 HOBO weather stations: Cavendish, Greenwich, North Rustico Wharf, Stanley Bridge Wharf, Tracadie Wharf",
            "ECCC Stanhope reference station (Climate ID 8300590) for cross-validation",
            "Agent-based architecture: Ingest > Clean > Redundancy > FWI > Uncertainty",
        ],
    },
    {
        "title": "Data Ingestion & Cleaning",
        "bullets": [
            "Ingested 5,086 daily observations across 20 variables",
            "100% row retention after cleaning (5,086 > 5,086 rows)",
            "Column reduction: 20 to 16 columns (dropped redundant/empty fields)",
            "2 temperature range issues flagged and handled",
            "Output: cleaned_data.parquet (380 KB)",
        ],
    },
    {
        "title": "Redundancy Analysis (PCA & Clustering)",
        "bullets": [
            "PCA: 4 components capture >90% of variance from 14 sensor channels",
            "10 channels carry redundant information - major optimization opportunity",
            "Wind/gust near-duplicate (r=0.96): decommission gust sensors at peripheral stations",
            "Pressure variables strongly correlated: shared synoptic weather across park",
            "3 K-Means clusters: Cluster 0 (351 extreme days) = all stations active",
            "Clusters 1-2 (~93% of days): reduced 3-station network sufficient",
        ],
        "image": "correlation_heatmap.png",
    },
    {
        "title": "Fire Weather Index (FWI)",
        "bullets": [
            "Van Wagner (1987) Canadian FWI System - 6 indices computed daily",
            "FWI range: 0.68 - 55.79 (mean 13.06)",
            "All 6 indices passed physical range validation (PASSED)",
            "Peak FWI = 55.8 at Stanley Bridge Wharf",
            "High FWI days (>=10): significant count; Very-High (>=20): flagged for review",
        ],
        "image": "fwi_plot.png",
    },
    {
        "title": "FWI Cross-Validation vs Stanhope",
        "bullets": [
            "Independent validation against ECCC Stanhope reference (368 fire-season days)",
            "FFMC MAE = 0.025 - near-identical",
            "DMC MAE = 0.70, DC MAE = 53.29 (cumulative index, expected divergence)",
            "ISI MAE = 0.015, BUI MAE = 3.20",
            "FWI MAE = 0.37 - excellent agreement",
        ],
    },
    {
        "title": "Uncertainty & Risk Assessment",
        "bullets": [
            "Total Variation (TV) distance: measures cleaning impact on data distributions",
            "HIGH RISK (>20% shift): water_level (47.9%), water_temperature (42.0%), pressures (36-38%)",
            "MODERATE RISK (5-20% shift): wind (12.7%), temperature (9.3%), wind_direction (7.4%)",
            "LOW RISK (<5% shift): rain (1.1%) - critical for FWI Drought Code, highly reliable",
            "Core FWI inputs (temp, humidity, wind, rain) all <10% shift - FWI calculations trustworthy",
            "Water sensors need field calibration before 2026 fire season (do not affect FWI directly)",
        ],
    },
    {
        "title": "Recommendations for Parks Canada",
        "bullets": [
            "1. DECOMMISSION gust-speed sensors at 2-3 stations (derive from wind, r=0.96, <4% error)",
            "2. CONSOLIDATE pressure monitoring to Cavendish (single central sensor replaces 5)",
            "3. MAINTAIN all 5 stations for temp, humidity, wind, rain (4 independent FWI inputs)",
            "4. SERVICE water_level and water_temperature sensors (42-48% distribution shift = drift)",
            "5. CONTINUE full network during fire season (May-Oct); reduced 3-station mode off-season",
            "Confidence: 5,086 days QC data, PCA, FWI cross-validation (MAE=0.37), KDE uncertainty",
        ],
    },
    {
        "title": "Repository & Reproducibility",
        "bullets": [
            "GitHub: github.com/Pegjoey99/parks-meteo-optimization",
            "Jupyter notebook: notebooks/analysis.ipynb - full end-to-end pipeline",
            "All outputs reproducible via: Run All in VS Code / Jupyter",
            "Dependencies: pip install -r requirements.txt",
            "License: MIT",
        ],
    },
]


# -- PowerPoint --------------------------------------------------------------

def build_pptx(output_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    DARK_BG = RGBColor(0x1B, 0x2A, 0x4A)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    ACCENT = RGBColor(0x4E, 0xC9, 0xB0)
    LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)

    def set_slide_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    # -- Title slide -----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = TITLE
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    p2 = tf.add_paragraph()
    p2.text = SUBTITLE
    p2.font.size = Pt(22)
    p2.font.color.rgb = ACCENT
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(12)

    p3 = tf.add_paragraph()
    p3.text = DATE
    p3.font.size = Pt(18)
    p3.font.color.rgb = LIGHT_GRAY
    p3.alignment = PP_ALIGN.LEFT
    p3.space_before = Pt(24)

    # -- Content slides --------------------------------------------------
    for s in SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        set_slide_bg(slide, DARK_BG)

        # Title bar
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = s["title"]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = ACCENT

        # Accent line
        line = slide.shapes.add_shape(
            1, Inches(0.8), Inches(1.35), Inches(11.5), Pt(2)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT
        line.line.fill.background()

        # Bullets
        has_image = "image" in s and (DATA_DIR / s["image"]).exists()
        bullet_width = Inches(5.5) if has_image else Inches(11.5)

        bullet_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.6), bullet_width, Inches(5.2)
        )
        tf = bullet_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(s["bullets"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = WHITE
            p.space_before = Pt(10)
            p.level = 0

        # Image (right side)
        if has_image:
            img_path = str(DATA_DIR / s["image"])
            slide.shapes.add_picture(
                img_path, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2)
            )

    prs.save(str(output_path))
    print(f"PowerPoint saved: {output_path} ({output_path.stat().st_size / 1024:.0f} KB)")


# -- PDF Report --------------------------------------------------------------

class PDFReport(FPDF):
    def header(self):
        self.set_fill_color(27, 42, 74)
        self.rect(0, 0, 210, 15, "F")
        self.set_font("Helvetica", "B", 10)
        self.set_text_color(78, 201, 176)
        self.set_y(4)
        self.cell(0, 7, TITLE, align="L")
        self.set_text_color(200, 200, 200)
        self.cell(0, 7, DATE, align="R")
        self.ln(12)

    def footer(self):
        self.set_y(-15)
        self.set_font("Helvetica", "I", 8)
        self.set_text_color(150, 150, 150)
        self.cell(0, 10, f"Page {self.page_no()}/{{nb}}", align="C")

    def section_title(self, title):
        self.set_font("Helvetica", "B", 16)
        self.set_text_color(27, 42, 74)
        self.cell(0, 12, title, new_x="LMARGIN", new_y="NEXT")
        self.set_draw_color(78, 201, 176)
        self.set_line_width(0.8)
        self.line(10, self.get_y(), 200, self.get_y())
        self.ln(4)

    def body_text(self, text):
        self.set_font("Helvetica", "", 11)
        self.set_text_color(40, 40, 40)
        self.multi_cell(0, 6, text)
        self.ln(2)

    def bullet_list(self, items):
        self.set_font("Helvetica", "", 11)
        self.set_text_color(40, 40, 40)
        for item in items:
            self.cell(5)
            self.cell(5, 6, "-")
            self.multi_cell(175, 6, item)
            self.ln(1)
        self.ln(3)

    def add_image_safe(self, path, w=170):
        if Path(path).exists():
            self.image(path, x=20, w=w)
            self.ln(5)


def build_pdf(output_path: Path):
    pdf = PDFReport()
    pdf.alias_nb_pages()
    pdf.set_auto_page_break(auto=True, margin=20)

    # -- Title page ------------------------------------------------------
    pdf.add_page()
    pdf.ln(50)
    pdf.set_font("Helvetica", "B", 32)
    pdf.set_text_color(27, 42, 74)
    pdf.cell(0, 15, TITLE, align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(5)
    pdf.set_font("Helvetica", "", 16)
    pdf.set_text_color(78, 201, 176)
    pdf.cell(0, 10, SUBTITLE, align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(10)
    pdf.set_font("Helvetica", "", 14)
    pdf.set_text_color(120, 120, 120)
    pdf.cell(0, 10, DATE, align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(20)
    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(80, 80, 80)
    pdf.cell(0, 8, "github.com/Pegjoey99/parks-meteo-optimization", align="C")

    # -- Content pages ---------------------------------------------------
    for s in SLIDES:
        pdf.add_page()
        pdf.section_title(s["title"])
        pdf.bullet_list(s["bullets"])

        if "image" in s:
            img_path = DATA_DIR / s["image"]
            pdf.add_image_safe(str(img_path), w=160)

    pdf.output(str(output_path))
    print(f"PDF saved: {output_path} ({output_path.stat().st_size / 1024:.0f} KB)")


# -- Main --------------------------------------------------------------------

if __name__ == "__main__":
    build_pptx(OUT_DIR / "PEINP_Parks_Meteo_Optimization.pptx")
    build_pdf(OUT_DIR / "PEINP_Parks_Meteo_Optimization.pdf")
    print("\nDone. Files ready in outputs/")
